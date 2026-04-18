from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ─── PALETTE ────────────────────────────────────────────────────────────────
BG       = RGBColor(0xF7, 0xF7, 0xF5)   # warm off-white
DARK     = RGBColor(0x1C, 0x1C, 0x1E)   # near-black
MID      = RGBColor(0x6E, 0x6E, 0x73)   # mid grey
ACCENT   = RGBColor(0x2C, 0x6E, 0x49)   # forest green
ACCENT2  = RGBColor(0x4A, 0xA8, 0x6C)   # lighter green
LINE     = RGBColor(0xD1, 0xD1, 0xD1)   # subtle divider

# ─── HELPERS ────────────────────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    return shape

def add_text(slide, text, l, t, w, h, size, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Helvetica Neue"
    return tb

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_line(slide, l, t, w, color=LINE, thickness=Pt(0.75)):
    ln = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(0.01))
    ln.fill.solid()
    ln.fill.fore_color.rgb = color
    ln.line.fill.background()
    return ln

# ─── PRESENTATION SETUP ─────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]   # completely blank

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 – PORTADA
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, BG)

# left accent bar
add_rect(s, 0, 0, 0.5, 7.5, ACCENT)

# title
add_text(s, "MindTracker", 1.1, 2.2, 9, 1.2, 52, bold=True, color=DARK)
add_text(s, "Aplicación web de seguimiento emocional y hábitos diarios",
         1.1, 3.5, 9, 0.7, 16, color=MID)

# divider
add_line(s, 1.1, 4.35, 6, color=ACCENT, thickness=Pt(1.5))

# meta
add_text(s, "Proyecto Integrado  ·  DAW 2º Curso  ·  2025-2026",
         1.1, 4.6, 9, 0.5, 11, color=MID)

# subtle dot decoration
add_rect(s, 10.8, 5.8, 0.18, 0.18, ACCENT2)
add_rect(s, 11.2, 5.8, 0.18, 0.18, LINE)
add_rect(s, 11.6, 5.8, 0.18, 0.18, LINE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 – ÍNDICE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, BG)
add_rect(s, 0, 0, 0.5, 7.5, ACCENT)
add_text(s, "Contenido", 1.1, 0.45, 10, 0.7, 28, bold=True, color=DARK)
add_line(s, 1.1, 1.25, 10.7, color=LINE)

items = [
    ("01", "Motivación y problema"),
    ("02", "Objetivos del proyecto"),
    ("03", "Stack tecnológico"),
    ("04", "Arquitectura del sistema"),
    ("05", "Base de datos"),
    ("06", "Backend · API REST"),
    ("07", "Frontend · React SPA"),
    ("08", "Funcionalidades clave"),
    ("09", "Seguridad"),
    ("10", "Despliegue"),
    ("11", "Pruebas"),
    ("12", "Conclusiones y trabajo futuro"),
]

cols = [items[:6], items[6:]]
for ci, col in enumerate(cols):
    x = 1.1 + ci * 6
    for i, (num, label) in enumerate(col):
        y = 1.5 + i * 0.85
        add_text(s, num, x, y, 0.7, 0.6, 11, bold=True, color=ACCENT)
        add_text(s, label, x + 0.7, y, 4.8, 0.6, 13, color=DARK)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 – MOTIVACIÓN
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, BG)
add_rect(s, 0, 0, 0.5, 7.5, ACCENT)
add_text(s, "01  Motivación", 1.1, 0.45, 10, 0.7, 22, bold=True, color=DARK)
add_line(s, 1.1, 1.25, 10.7, color=LINE)

add_text(s, "El problema", 1.1, 1.5, 5, 0.5, 15, bold=True, color=DARK)
add_text(s,
    "Las apps comerciales de bienestar (Daylio, Bearable…) exigen suscripciones de pago, "
    "almacenan datos personales sensibles en la nube de terceros y limitan la personalización "
    "del usuario. No existe una alternativa libre, autohospedada y de código abierto con una "
    "experiencia de usuario comparable.",
    1.1, 2.05, 11, 1.4, 13, color=DARK)

add_text(s, "La oportunidad", 1.1, 3.7, 5, 0.5, 15, bold=True, color=DARK)
add_text(s,
    "Construir una herramienta propia que devolviese el control de los datos al usuario, "
    "utilizando tecnologías del currículo DAW — Node.js, React, MariaDB y Apache — "
    "en un entorno de despliegue real sobre Ubuntu 22.04.",
    1.1, 4.25, 11, 1.3, 13, color=DARK)

# accent quote block
add_rect(s, 1.1, 5.8, 0.08, 0.9, ACCENT)
add_text(s, "«Cada persona debería ser dueña de sus propios datos de bienestar.»",
         1.35, 5.85, 10.5, 0.7, 12, italic=True, color=MID)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 – OBJETIVOS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, BG)
add_rect(s, 0, 0, 0.5, 7.5, ACCENT)
add_text(s, "02  Objetivos", 1.1, 0.45, 10, 0.7, 22, bold=True, color=DARK)
add_line(s, 1.1, 1.25, 10.7, color=LINE)

gen_obj = "Desarrollar una aplicación web full-stack funcional, segura y desplegable en producción para el registro y seguimiento de estados emocionales y hábitos diarios."
add_text(s, "Objetivo general", 1.1, 1.5, 11, 0.4, 13, bold=True, color=ACCENT)
add_text(s, gen_obj, 1.1, 1.95, 11, 0.8, 12, color=DARK)
add_line(s, 1.1, 2.85, 10.7, color=LINE)

specifics = [
    "Implementar autenticación segura con JWT en cookies HttpOnly y contraseñas hasheadas con bcrypt.",
    "Diseñar e implementar una API REST completa con Express.js sobre una base de datos MariaDB normalizada.",
    "Desarrollar una SPA con React 18 que ofrezca una UX fluida con registro de estado de ánimo y hábitos.",
    "Desplegar el sistema completo en Ubuntu 22.04 con Apache como proxy inverso y SSL/TLS.",
    "Cubrir los 5 módulos del ciclo (Servidor, Cliente, BBDD, Despliegue, Sistemas) en un proyecto real.",
]
add_text(s, "Objetivos específicos", 1.1, 3.05, 11, 0.4, 13, bold=True, color=DARK)
for i, sp in enumerate(specifics):
    y = 3.5 + i * 0.62
    add_rect(s, 1.1, y + 0.12, 0.12, 0.12, ACCENT)
    add_text(s, sp, 1.35, y, 11, 0.55, 11.5, color=DARK)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 – STACK TECNOLÓGICO
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, BG)
add_rect(s, 0, 0, 0.5, 7.5, ACCENT)
add_text(s, "03  Stack tecnológico", 1.1, 0.45, 10, 0.7, 22, bold=True, color=DARK)
add_line(s, 1.1, 1.25, 10.7, color=LINE)

categories = [
    ("Frontend",  ["React 18", "React Router v6", "Vite", "CSS + FontAwesome"]),
    ("Backend",   ["Node.js 20", "Express.js 5", "JWT (jsonwebtoken)", "bcrypt"]),
    ("Base de datos", ["MariaDB 10.6", "mysql2 connector", "Pool de conexiones"]),
    ("Infraestructura", ["Apache 2.4", "Ubuntu 22.04 (WSL2)", "mod_ssl / HTTPS", "systemd"]),
]

for ci, (cat, techs) in enumerate(categories):
    col = ci % 2
    row = ci // 2
    x = 1.1 + col * 6
    y = 1.6 + row * 2.8

    # card bg
    card = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(5.5), Inches(2.4))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xED, 0xED, 0xEB)
    card.line.fill.background()

    add_rect(s, x, y, 5.5, 0.06, ACCENT)
    add_text(s, cat, x + 0.2, y + 0.15, 5, 0.45, 13, bold=True, color=DARK)
    for ti, tech in enumerate(techs):
        add_rect(s, x + 0.2, y + 0.75 + ti * 0.38, 0.07, 0.07, ACCENT2)
        add_text(s, tech, x + 0.4, y + 0.7 + ti * 0.38, 4.8, 0.38, 11.5, color=DARK)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 – ARQUITECTURA
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, BG)
add_rect(s, 0, 0, 0.5, 7.5, ACCENT)
add_text(s, "04  Arquitectura del sistema", 1.1, 0.45, 10, 0.7, 22, bold=True, color=DARK)
add_line(s, 1.1, 1.25, 10.7, color=LINE)

layers = [
    ("Capa de Presentación",   "React SPA · client-side routing · componentes reutilizables"),
    ("Proxy Inverso",          "Apache 2.4 · rutas /api/* → Node.js · archivos estáticos"),
    ("Capa de Aplicación",     "Express.js API · MVC · Middleware JWT · CORS"),
    ("Capa de Persistencia",   "MariaDB · 5 tablas normalizadas · foreign keys · connection pool"),
]
colors_l = [ACCENT, ACCENT2, RGBColor(0x74,0xB6,0x8A), RGBColor(0xA2,0xCC,0xB0)]

for i, (title, desc) in enumerate(layers):
    y = 1.65 + i * 1.3
    add_rect(s, 1.1, y, 11, 1.0, colors_l[i])
    add_text(s, title, 1.3, y + 0.08, 5, 0.45, 12, bold=True, color=RGBColor(0xFF,0xFF,0xFF))
    add_text(s, desc,  1.3, y + 0.5,  10.5, 0.45, 11, color=RGBColor(0xF0,0xF0,0xF0))
    if i < 3:
        add_text(s, "↓", 6.45, y + 1.0, 0.5, 0.4, 14, bold=True, color=MID, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 – BASE DE DATOS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, BG)
add_rect(s, 0, 0, 0.5, 7.5, ACCENT)
add_text(s, "05  Base de datos", 1.1, 0.45, 10, 0.7, 22, bold=True, color=DARK)
add_line(s, 1.1, 1.25, 10.7, color=LINE)

tables = [
    ("users",        ["id (PK)", "name, last_name", "email (UNIQUE)", "phone", "password (bcrypt)"]),
    ("moods",        ["id (PK)", "user_id (FK→users)", "value 1-10", "notes", "date"]),
    ("habits",       ["id (PK)", "name", "icon"]),
    ("habit_options",["id (PK)", "habit_id (FK)", "label", "sort_order"]),
    ("habit_logs",   ["id (PK)", "user_id (FK)", "habit_id (FK)", "option_id (FK)", "date · UNIQUE(user,habit,date)"]),
]

col_widths = [2.2, 2.2, 2.2, 2.5, 3.4]
xs = [1.1, 3.4, 5.6, 7.85, 10.4]

for ci, ((tname, fields), x, w) in enumerate(zip(tables, xs, col_widths)):
    y = 1.55
    add_rect(s, x, y, w, 0.48, ACCENT)
    add_text(s, tname, x + 0.1, y + 0.05, w - 0.1, 0.38, 11, bold=True,
             color=RGBColor(0xFF,0xFF,0xFF))
    for fi, field in enumerate(fields):
        fy = y + 0.48 + fi * 0.5
        bg_col = RGBColor(0xED,0xED,0xEB) if fi % 2 == 0 else BG
        add_rect(s, x, fy, w, 0.5, bg_col)
        add_text(s, field, x + 0.1, fy + 0.08, w - 0.1, 0.38, 9.5, color=DARK)

add_text(s, "Restricciones: FK con ON DELETE CASCADE · UNIQUE(user_id, habit_id, date) en habit_logs · email UNIQUE en users",
         1.1, 6.6, 11.5, 0.5, 10, italic=True, color=MID)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 – BACKEND
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, BG)
add_rect(s, 0, 0, 0.5, 7.5, ACCENT)
add_text(s, "06  Backend · API REST", 1.1, 0.45, 10, 0.7, 22, bold=True, color=DARK)
add_line(s, 1.1, 1.25, 10.7, color=LINE)

# left: endpoints table
add_text(s, "Endpoints principales", 1.1, 1.5, 6, 0.4, 13, bold=True, color=DARK)
endpoints = [
    ("POST", "/auth/register",          "Registro de usuario"),
    ("POST", "/auth/login",             "Login · emite JWT en cookie"),
    ("POST", "/auth/logout",            "Logout · limpia cookie"),
    ("GET",  "/api/me",                 "Usuario autenticado actual"),
    ("POST", "/api/users/moods",        "Crear entrada de estado de ánimo"),
    ("GET",  "/api/users/logs",         "Historial de ánimo + hábitos"),
    ("PUT",  "/api/users/moods/:id",    "Actualizar entrada"),
    ("DELETE","/api/users/moods/:id",   "Eliminar entrada"),
    ("GET",  "/api/habits",             "Catálogo de hábitos con opciones"),
    ("POST", "/api/habits/logs",        "Registrar hábito (upsert)"),
    ("GET",  "/api/habits/logs/date/:d","Hábitos de una fecha específica"),
]
method_colors = {"GET": RGBColor(0x2C,0x6E,0x49), "POST": RGBColor(0x27,0x5D,0x8E),
                 "PUT": RGBColor(0xA0,0x6E,0x1A), "DELETE": RGBColor(0x8E,0x27,0x27)}
for i, (method, path, desc) in enumerate(endpoints):
    y = 2.0 + i * 0.43
    mc = method_colors.get(method, DARK)
    add_text(s, method, 1.1, y, 0.85, 0.38, 8.5, bold=True, color=mc)
    add_text(s, path,   1.95, y, 3.2,  0.38, 8.5, color=DARK)
    add_text(s, desc,   5.2,  y, 2.8,  0.38, 8.5, color=MID)

# right: MVC pattern
add_text(s, "Patrón MVC", 8.5, 1.5, 4, 0.4, 13, bold=True, color=DARK)
mvc = [
    ("Routes",      "Definición de endpoints y vinculación con controladores"),
    ("Controllers", "Lógica de negocio, validaciones, respuesta HTTP"),
    ("Models",      "Queries parametrizadas a MariaDB · pool de conexiones"),
    ("Middleware",  "Verificación JWT · protección de rutas privadas"),
]
for i, (comp, desc) in enumerate(mvc):
    y = 2.0 + i * 1.1
    add_rect(s, 8.5, y, 4.2, 0.95, RGBColor(0xED,0xED,0xEB))
    add_rect(s, 8.5, y, 0.08, 0.95, ACCENT)
    add_text(s, comp, 8.7, y + 0.05, 3.8, 0.38, 11, bold=True, color=DARK)
    add_text(s, desc, 8.7, y + 0.46, 3.8, 0.42, 9.5, color=MID)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 – FRONTEND
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, BG)
add_rect(s, 0, 0, 0.5, 7.5, ACCENT)
add_text(s, "07  Frontend · React SPA", 1.1, 0.45, 10, 0.7, 22, bold=True, color=DARK)
add_line(s, 1.1, 1.25, 10.7, color=LINE)

# component tree
add_text(s, "Árbol de componentes", 1.1, 1.5, 8, 0.4, 13, bold=True, color=DARK)
tree = [
    (0, "App.jsx",           "Routing principal · ProtectedRoute / PublicRoute"),
    (1, "Login.jsx",         "Formularios de acceso y registro con validación"),
    (1, "Dashboard.jsx",     "Shell principal con AppHeader y vistas dinámicas"),
    (2, "EntradaView",       "Formulario de estado de ánimo + panel de hábitos"),
    (3, "MoodForm",          "Grid 1–10 · notas · selector de fecha"),
    (3, "HabitsPanel",       "5 hábitos · opciones dinámicas · guardado inmediato"),
    (2, "RegistrosView",     "Historial cronológico · eliminación · tags de hábitos"),
    (2, "PerfilView",        "Edición de perfil · cambio de contraseña"),
]
for i, (level, comp, desc) in enumerate(tree):
    y = 1.95 + i * 0.57
    indent = 1.1 + level * 0.45
    if level > 0:
        add_text(s, "└─", indent - 0.38, y, 0.38, 0.45, 9, color=LINE)
    add_text(s, comp, indent, y, 2.8, 0.45, 10.5, bold=(level < 2), color=DARK)
    add_text(s, desc, indent + 2.85, y, 6.5, 0.45, 10, color=MID)

add_line(s, 1.1, 6.5, 10.7, color=LINE)
decisions = "SPA con React Router · Estado local (useState/useEffect) · Fetch API con credentials · Rutas protegidas por cookie HttpOnly"
add_text(s, "Decisiones de diseño:  " + decisions, 1.1, 6.55, 11.5, 0.6, 10, italic=True, color=MID)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 – FUNCIONALIDADES
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, BG)
add_rect(s, 0, 0, 0.5, 7.5, ACCENT)
add_text(s, "08  Funcionalidades clave", 1.1, 0.45, 10, 0.7, 22, bold=True, color=DARK)
add_line(s, 1.1, 1.25, 10.7, color=LINE)

features = [
    ("Estado de ánimo",
     ["Escala 1-10 como grid de botones", "Notas de texto libre", "Selector de fecha",
      "CRUD completo sobre cada entrada"]),
    ("Seguimiento de hábitos",
     ["5 hábitos predefinidos (sueño, ejercicio,\nlectura, estudio, dieta)",
      "Opciones múltiples por hábito",
      "Guardado inmediato (sin botón extra)",
      "Upsert: un registro por hábito/día/usuario"]),
    ("Historial y perfil",
     ["Vista cronológica con tags de hábitos", "Eliminación de entradas",
      "Edición de nombre, apellidos y teléfono",
      "Cambio de contraseña con verificación actual"]),
]

for ci, (title, points) in enumerate(features):
    x = 1.1 + ci * 4.1
    y = 1.6
    card = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(3.9), Inches(5.4))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xED,0xED,0xEB)
    card.line.fill.background()
    add_rect(s, x, y, 3.9, 0.06, ACCENT)
    add_text(s, title, x + 0.2, y + 0.12, 3.5, 0.5, 13, bold=True, color=DARK)
    add_line(s, x + 0.2, y + 0.7, 3.5, color=LINE)
    for pi, pt in enumerate(points):
        py = y + 0.9 + pi * 1.05
        add_rect(s, x + 0.2, py + 0.1, 0.1, 0.1, ACCENT2)
        add_text(s, pt, x + 0.42, py, 3.3, 0.95, 11, color=DARK)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 – SEGURIDAD
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, BG)
add_rect(s, 0, 0, 0.5, 7.5, ACCENT)
add_text(s, "09  Seguridad", 1.1, 0.45, 10, 0.7, 22, bold=True, color=DARK)
add_line(s, 1.1, 1.25, 10.7, color=LINE)

security_items = [
    ("JWT en cookies HttpOnly",
     "El token no es accesible desde JavaScript, eliminando el riesgo de robo mediante XSS."),
    ("Hashing con bcrypt (10 rounds)",
     "Las contraseñas nunca se almacenan en texto plano. El factor de coste protege contra fuerza bruta."),
    ("Queries parametrizadas",
     "Todas las consultas a MariaDB usan placeholders — inmunidad total a SQL Injection."),
    ("Cabeceras HTTP de seguridad",
     "Apache añade X-Frame-Options, X-Content-Type-Options y X-XSS-Protection a todas las respuestas."),
    ("CORS restringido",
     "El backend solo acepta peticiones del origen configurado, bloqueando CSRF cross-origin."),
    ("SSL/TLS en producción",
     "Comunicación cifrada end-to-end con certificado servido desde Apache mod_ssl."),
]

for i, (title, desc) in enumerate(security_items):
    col = i % 2
    row = i // 2
    x = 1.1 + col * 6
    y = 1.6 + row * 1.7
    card = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(5.6), Inches(1.5))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xED,0xED,0xEB)
    card.line.fill.background()
    add_rect(s, x, y, 0.08, 1.5, ACCENT)
    add_text(s, title, x + 0.22, y + 0.1, 5.2, 0.45, 12, bold=True, color=DARK)
    add_text(s, desc,  x + 0.22, y + 0.6, 5.2, 0.82, 10.5, color=MID)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 – DESPLIEGUE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, BG)
add_rect(s, 0, 0, 0.5, 7.5, ACCENT)
add_text(s, "10  Despliegue", 1.1, 0.45, 10, 0.7, 22, bold=True, color=DARK)
add_line(s, 1.1, 1.25, 10.7, color=LINE)

steps = [
    ("1", "Ubuntu 22.04 (WSL2)", "Entorno de servidor real con acceso root y gestión de servicios con systemd"),
    ("2", "MariaDB 10.6",        "Base de datos instalada y configurada con usuario dedicado y schema MindTracker"),
    ("3", "Node.js 20 + API",    "Servidor Express corriendo como servicio con PM2/systemd en el puerto 2345"),
    ("4", "Build React (Vite)",  "Compilación del frontend con 'vite build' generando el directorio dist/"),
    ("5", "Apache + SSL",        "Proxy inverso con VirtualHost: /api/* → Node.js, /* → archivos estáticos"),
    ("6", "Variables de entorno",".env con JWT_SECRET, DB_PASS, ORIGIN — nunca en control de versiones"),
]

for i, (num, title, desc) in enumerate(steps):
    col = i % 2
    row = i // 2
    x = 1.1 + col * 6
    y = 1.6 + row * 1.7
    add_rect(s, x, y + 0.08, 0.5, 0.5, ACCENT)
    add_text(s, num, x, y + 0.08, 0.5, 0.5, 16, bold=True,
             color=RGBColor(0xFF,0xFF,0xFF), align=PP_ALIGN.CENTER)
    add_text(s, title, x + 0.65, y + 0.08, 5, 0.45, 12, bold=True, color=DARK)
    add_text(s, desc,  x + 0.65, y + 0.55, 5, 0.9,  10.5, color=MID)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 – PRUEBAS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, BG)
add_rect(s, 0, 0, 0.5, 7.5, ACCENT)
add_text(s, "11  Pruebas funcionales", 1.1, 0.45, 10, 0.7, 22, bold=True, color=DARK)
add_line(s, 1.1, 1.25, 10.7, color=LINE)

add_text(s, "14 casos de prueba documentados · todos superados ✓", 1.1, 1.5, 11, 0.5, 13,
         bold=True, color=ACCENT)

test_cases = [
    ("CP-01", "Registro de nuevo usuario con datos válidos"),
    ("CP-02", "Registro con email duplicado — error esperado"),
    ("CP-03", "Login con credenciales correctas"),
    ("CP-04", "Login con contraseña incorrecta — error esperado"),
    ("CP-05", "Acceso a ruta protegida sin sesión — redirección"),
    ("CP-06", "Registro de estado de ánimo en fecha actual"),
    ("CP-07", "Visualización del historial de ánimo"),
    ("CP-08", "Eliminación de una entrada de historial"),
    ("CP-09", "Registro de hábito · guardado inmediato"),
    ("CP-10", "Actualización de hábito ya registrado (upsert)"),
    ("CP-11", "Visualización de hábitos del día actual"),
    ("CP-12", "Actualización de datos de perfil"),
    ("CP-13", "Cambio de contraseña con verificación correcta"),
    ("CP-14", "Logout · invalidación de sesión"),
]

cols2 = [test_cases[:7], test_cases[7:]]
for ci, col in enumerate(cols2):
    x = 1.1 + ci * 6
    for i, (code, desc) in enumerate(col):
        y = 2.2 + i * 0.6
        add_text(s, code, x, y, 0.9, 0.52, 9.5, bold=True, color=ACCENT)
        add_text(s, desc, x + 0.9, y, 5, 0.52, 10.5, color=DARK)
        add_text(s, "✓", x + 5.9, y, 0.4, 0.52, 10, bold=True,
                 color=ACCENT, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 – CONCLUSIONES
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, BG)
add_rect(s, 0, 0, 0.5, 7.5, ACCENT)
add_text(s, "12  Conclusiones y trabajo futuro", 1.1, 0.45, 11, 0.7, 22, bold=True, color=DARK)
add_line(s, 1.1, 1.25, 10.7, color=LINE)

conclusions = [
    "Se ha desarrollado y desplegado con éxito una aplicación web full-stack funcional y segura.",
    "Los 5 módulos del ciclo DAW quedan integrados de forma coherente en un único proyecto real.",
    "El sistema implementa las mejores prácticas de seguridad web: JWT, bcrypt, parameterized queries y HTTPS.",
    "La arquitectura en tres capas facilita el mantenimiento y la escalabilidad futura.",
    "Los 14 casos de prueba validados demuestran la robustez y corrección de todas las funcionalidades.",
]
add_text(s, "Logros", 1.1, 1.5, 8, 0.4, 13, bold=True, color=DARK)
for i, c in enumerate(conclusions):
    y = 1.95 + i * 0.55
    add_rect(s, 1.1, y + 0.1, 0.1, 0.1, ACCENT)
    add_text(s, c, 1.32, y, 11.2, 0.5, 11, color=DARK)

add_line(s, 1.1, 4.9, 10.7, color=LINE)
add_text(s, "Trabajo futuro", 1.1, 5.05, 8, 0.4, 13, bold=True, color=DARK)
future = [
    ("Gráficas y tendencias",   "Visualización de evolución del estado de ánimo con Chart.js"),
    ("Hábitos personalizables", "CRUD de hábitos y opciones definidos por el usuario"),
    ("Aplicación móvil",        "PWA o app nativa con React Native para registro rápido"),
    ("Notificaciones",          "Recordatorios diarios vía email o push notification"),
]
for ci, (title, desc) in enumerate(future):
    col = ci % 2
    row = ci // 2
    x = 1.1 + col * 6
    y = 5.55 + row * 0.62
    add_rect(s, x, y + 0.1, 0.1, 0.1, ACCENT2)
    add_text(s, title + ": ", x + 0.25, y, 1.8, 0.52, 10.5, bold=True, color=DARK)
    add_text(s, desc, x + 2.1, y, 3.7, 0.52, 10.5, color=MID)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 – CIERRE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, BG)
add_rect(s, 0, 0, 0.5, 7.5, ACCENT)
add_rect(s, 0.5, 0, 12.83, 7.5, BG)  # ensure clean bg

add_text(s, "Gracias", 1.1, 2.2, 11, 1.2, 60, bold=True, color=DARK, align=PP_ALIGN.CENTER)
add_line(s, 3.5, 3.6, 6.3, color=ACCENT, thickness=Pt(1.5))
add_text(s, "MindTracker · Proyecto Integrado DAW · 2025-2026",
         1.1, 3.85, 11, 0.6, 13, color=MID, align=PP_ALIGN.CENTER)
add_text(s, "¿Preguntas?",
         1.1, 4.6, 11, 0.6, 16, color=DARK, align=PP_ALIGN.CENTER)

# ─── SAVE ────────────────────────────────────────────────────────────────────
output = "/home/mbalay19/mindtracker-copy/MindTracker_Presentacion.pptx"
prs.save(output)
print(f"Saved: {output}")
